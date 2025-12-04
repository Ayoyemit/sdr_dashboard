import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Helper to save plotly fig as image
import plotly.io as pio

def create_sample_sankey(title):
    fig = go.Figure(go.Sankey(
        node=dict(
            pad=15,
            thickness=20,
            line=dict(color="black", width=0.5),
            label=["Mothers", "Low Risk", "High Risk", "L4", "L5"],
            color=["#4169e1", "#90ee90", "#ff8c00", "#ffd700", "#8a2be2"]
        ),
        link=dict(
            source=[0, 0, 1, 2],
            target=[1, 2, 3, 4],
            value=[60, 40, 30, 20],
            color=["rgba(65,105,225,0.2)", "rgba(255,140,0,0.2)", "rgba(255,215,0,0.2)", "rgba(138,43,226,0.2)"]
        )
    ))
    fig.update_layout(title_text=title, font_size=12, width=700, height=400)
    return fig

def create_architecture_visual(path):
    # Use plotly to create a simple architecture diagram
    fig = go.Figure()
    # Dashboard
    fig.add_shape(type="rect", x0=0.1, y0=0.7, x1=0.3, y1=0.9, line=dict(color="RoyalBlue"), fillcolor="#e6f0ff")
    fig.add_annotation(x=0.2, y=0.8, text="Dashboard UI", showarrow=False, font=dict(size=14))
    # Data sources
    fig.add_shape(type="rect", x0=0.1, y0=0.1, x1=0.3, y1=0.3, line=dict(color="green"), fillcolor="#eaffea")
    fig.add_annotation(x=0.2, y=0.2, text="Distributed Data Sources", showarrow=False, font=dict(size=12))
    # ETL
    fig.add_shape(type="rect", x0=0.4, y0=0.1, x1=0.6, y1=0.3, line=dict(color="orange"), fillcolor="#fff5e6")
    fig.add_annotation(x=0.5, y=0.2, text="ETL & Cleaning", showarrow=False, font=dict(size=12))
    # Backend
    fig.add_shape(type="rect", x0=0.7, y0=0.1, x1=0.9, y1=0.3, line=dict(color="purple"), fillcolor="#f3e6ff")
    fig.add_annotation(x=0.8, y=0.2, text="Distributed Compute/Cloud", showarrow=False, font=dict(size=12))
    # Arrows
    fig.add_annotation(x=0.35, y=0.2, ax=0.3, ay=0.2, xref='x', yref='y', axref='x', ayref='y', text='', showarrow=True, arrowhead=2)
    fig.add_annotation(x=0.65, y=0.2, ax=0.6, ay=0.2, xref='x', yref='y', axref='x', ayref='y', text='', showarrow=True, arrowhead=2)
    fig.add_annotation(x=0.2, y=0.7, ax=0.2, ay=0.3, xref='x', yref='y', axref='x', ayref='y', text='', showarrow=True, arrowhead=2)
    fig.update_layout(width=700, height=400, xaxis=dict(visible=False), yaxis=dict(visible=False), plot_bgcolor='white', margin=dict(l=0, r=0, t=0, b=0))
    pio.write_image(fig, path)

def add_slide(prs, title, bullet_points, image_path=None):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(3.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    if image_path:
        slide.shapes.add_picture(image_path, Inches(1), Inches(3.5), width=Inches(7))

def main():
    prs = Presentation()
    # Slide 1: Sankey Diagrams for Pathway Analysis
    sankey1 = create_sample_sankey("Risk Stratification & Delivery Location Pathway")
    sankey1_path = "sankey1.png"
    pio.write_image(sankey1, sankey1_path)
    add_slide(
        prs,
        "Sankey Diagrams for Pathway Analysis",
        [
            "• Sankey diagrams visualize patient journeys and intervention impacts across health pathways.",
            "• Pathways modeled: Risk Stratification, ANC & Complications, Intrapartum Monitoring, SDR Delivery.",
            "• Each stage is mapped to a node; flows represent transitions; color-coding highlights bottlenecks and effects."
        ],
        sankey1_path
    )
    # Slide 2: Comparative Insights
    sankey2 = create_sample_sankey("Risk Stratification & Delivery Location Pathway (Intervention)")
    sankey2_path = "sankey2.png"
    pio.write_image(sankey2, sankey2_path)
    add_slide(
        prs,
        "Comparative Insights from Sankey Visualizations",
        [
            "• Baseline vs. intervention Sankey diagrams reveal how interventions shift patient flows and outcomes.",
            "• Quantifies reductions in adverse events and improvements in care pathways.",
            "• Users can select pathways and scenarios in the dashboard for instant visual feedback."
        ],
        sankey2_path
    )
    # Slide 3: Deployment, Optimization, and Future Integration
    arch_path = "architecture.png"
    create_architecture_visual(arch_path)
    add_slide(
        prs,
        "Deployment, Optimization, and Future Integration",
        [
            "• Dashboard optimized for online use (Streamlit, Plotly, caching, memory management).",
            "• Future: Integrate with real-time, distributed health data (with missing/inconsistent data).",
            "• Requires robust ETL, scalable backend, advanced analytics, and compliance with health data regulations."
        ],
        arch_path
    )
    prs.save("Sankey_Dashboard_Presentation.pptx")
    # Clean up images
    os.remove(sankey1_path)
    os.remove(sankey2_path)
    os.remove(arch_path)

if __name__ == "__main__":
    main() 